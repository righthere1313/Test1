from __future__ import annotations

from typing import Dict, Optional, Tuple
import xml.etree.ElementTree as ET

from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Pt
from pptx.dml.color import RGBColor


def _tag_name(el: ET.Element) -> str:
    tag = el.tag or ""
    if "}" in tag:
        return tag.split("}", 1)[1]
    return tag


def _parse_viewbox(root: ET.Element) -> Tuple[float, float]:
    vb = (root.get("viewBox") or "").strip()
    if vb:
        parts = vb.replace(",", " ").split()
        if len(parts) >= 4:
            try:
                return float(parts[2]), float(parts[3])
            except Exception:
                pass
    try:
        w = float((root.get("width") or "1280").replace("px", "").strip())
        h = float((root.get("height") or "720").replace("px", "").strip())
        return w, h
    except Exception:
        return 1280.0, 720.0


def _float_attr(el: ET.Element, key: str, default: float = 0.0) -> float:
    raw = (el.get(key) or "").strip()
    if not raw:
        return default
    raw = raw.replace("px", "")
    try:
        return float(raw)
    except Exception:
        return default


def _hex_color(s: str) -> Optional[RGBColor]:
    raw = (s or "").strip()
    if not raw:
        return None
    if raw.startswith("#"):
        raw = raw[1:]
    if len(raw) == 3:
        raw = "".join([c * 2 for c in raw])
    if len(raw) != 6:
        return None
    try:
        r = int(raw[0:2], 16)
        g = int(raw[2:4], 16)
        b = int(raw[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return None


def _clamp01(x: float) -> float:
    try:
        return max(0.0, min(1.0, float(x)))
    except Exception:
        return 1.0


def _parse_opacity(el: ET.Element, key: str) -> Optional[float]:
    raw = (el.get(key) or "").strip()
    if not raw:
        return None
    try:
        return _clamp01(float(raw))
    except Exception:
        return None


def _parse_style_kv(style: str) -> Dict[str, str]:
    out: Dict[str, str] = {}
    raw = (style or "").strip()
    if not raw:
        return out
    for part in raw.split(";"):
        if ":" not in part:
            continue
        k, v = part.split(":", 1)
        k = k.strip()
        v = v.strip()
        if k:
            out[k] = v
    return out


def _resolve_paint(paint: str, gradients: Dict[str, str]) -> Optional[str]:
    raw = (paint or "").strip()
    if not raw or raw.lower() in {"none", "transparent"}:
        return None
    if raw.startswith("url(") and "#" in raw:
        gid = raw.split("#", 1)[1].split(")", 1)[0].strip()
        return gradients.get(gid)
    return raw


def _collect_linear_gradients(root: ET.Element) -> Dict[str, str]:
    grads: Dict[str, str] = {}
    for el in root.iter():
        if _tag_name(el) != "linearGradient":
            continue
        gid = (el.get("id") or "").strip()
        if not gid:
            continue
        first_color = None
        for st in el:
            if _tag_name(st) != "stop":
                continue
            style = _parse_style_kv(st.get("style") or "")
            c = (style.get("stop-color") or st.get("stop-color") or "").strip()
            if c:
                first_color = c
                break
        if first_color:
            grads[gid] = first_color
    return grads


def draw_svg_visuals(slide, prs, svg_text: str) -> None:
    raw = (svg_text or "").strip()
    if not raw:
        return
    try:
        root = ET.fromstring(raw.encode("utf-8", errors="ignore"))
    except Exception:
        return
    vw, vh = _parse_viewbox(root)
    vw = max(1.0, float(vw))
    vh = max(1.0, float(vh))
    sw = int(getattr(prs, "slide_width", 1) or 1)
    sh = int(getattr(prs, "slide_height", 1) or 1)
    gradients = _collect_linear_gradients(root)

    for el in list(root):
        tag = _tag_name(el)
        if tag == "defs":
            continue
        if tag == "text":
            continue
        if tag == "rect":
            x = _float_attr(el, "x", 0.0)
            y = _float_attr(el, "y", 0.0)
            w = _float_attr(el, "width", 0.0)
            h = _float_attr(el, "height", 0.0)
            if w <= 0 or h <= 0:
                continue
            left = int(x / vw * sw)
            top = int(y / vh * sh)
            width = int(w / vw * sw)
            height = int(h / vh * sh)
            rx = _float_attr(el, "rx", 0.0)
            shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rx > 0 else MSO_SHAPE.RECTANGLE
            try:
                shp = slide.shapes.add_shape(shape_type, left, top, width, height)
            except Exception:
                continue
            style_kv = _parse_style_kv(el.get("style") or "")
            fill = _resolve_paint(style_kv.get("fill") or el.get("fill") or "", gradients)
            stroke = _resolve_paint(style_kv.get("stroke") or el.get("stroke") or "", gradients)
            fill_op = _parse_opacity(el, "fill-opacity")
            stroke_op = _parse_opacity(el, "stroke-opacity")
            if fill:
                try:
                    shp.fill.solid()
                    c = _hex_color(fill)
                    if c is not None:
                        shp.fill.fore_color.rgb = c
                    if fill_op is not None:
                        shp.fill.fore_color.transparency = 1.0 - float(fill_op)
                except Exception:
                    pass
            else:
                try:
                    shp.fill.background()
                except Exception:
                    pass
            if stroke:
                try:
                    c = _hex_color(stroke)
                    if c is not None:
                        shp.line.color.rgb = c
                    if stroke_op is not None:
                        shp.line.transparency = 1.0 - float(stroke_op)
                    sw_px = _float_attr(el, "stroke-width", _float_attr(el, "strokeWidth", 1.0))
                    if sw_px > 0:
                        shp.line.width = Pt(sw_px * 0.75)
                except Exception:
                    pass
            else:
                try:
                    shp.line.fill.background()
                except Exception:
                    pass
            continue
        if tag == "circle":
            cx = _float_attr(el, "cx", 0.0)
            cy = _float_attr(el, "cy", 0.0)
            r = _float_attr(el, "r", 0.0)
            if r <= 0:
                continue
            x = cx - r
            y = cy - r
            w = r * 2.0
            h = r * 2.0
            left = int(x / vw * sw)
            top = int(y / vh * sh)
            width = int(w / vw * sw)
            height = int(h / vh * sh)
            try:
                shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
            except Exception:
                continue
            style_kv = _parse_style_kv(el.get("style") or "")
            fill = _resolve_paint(style_kv.get("fill") or el.get("fill") or "", gradients)
            stroke = _resolve_paint(style_kv.get("stroke") or el.get("stroke") or "", gradients)
            fill_op = _parse_opacity(el, "fill-opacity")
            stroke_op = _parse_opacity(el, "stroke-opacity")
            if fill:
                try:
                    shp.fill.solid()
                    c = _hex_color(fill)
                    if c is not None:
                        shp.fill.fore_color.rgb = c
                    if fill_op is not None:
                        shp.fill.fore_color.transparency = 1.0 - float(fill_op)
                except Exception:
                    pass
            else:
                try:
                    shp.fill.background()
                except Exception:
                    pass
            if stroke:
                try:
                    c = _hex_color(stroke)
                    if c is not None:
                        shp.line.color.rgb = c
                    if stroke_op is not None:
                        shp.line.transparency = 1.0 - float(stroke_op)
                except Exception:
                    pass
            else:
                try:
                    shp.line.fill.background()
                except Exception:
                    pass
            continue
        if tag == "line":
            x1 = _float_attr(el, "x1", 0.0)
            y1 = _float_attr(el, "y1", 0.0)
            x2 = _float_attr(el, "x2", 0.0)
            y2 = _float_attr(el, "y2", 0.0)
            start_x = int(x1 / vw * sw)
            start_y = int(y1 / vh * sh)
            end_x = int(x2 / vw * sw)
            end_y = int(y2 / vh * sh)
            try:
                conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y)
            except Exception:
                continue
            style_kv = _parse_style_kv(el.get("style") or "")
            stroke = _resolve_paint(style_kv.get("stroke") or el.get("stroke") or "", gradients)
            stroke_op = _parse_opacity(el, "stroke-opacity")
            if stroke:
                try:
                    c = _hex_color(stroke)
                    if c is not None:
                        conn.line.color.rgb = c
                    if stroke_op is not None:
                        conn.line.transparency = 1.0 - float(stroke_op)
                    sw_px = _float_attr(el, "stroke-width", _float_attr(el, "strokeWidth", 1.0))
                    if sw_px > 0:
                        conn.line.width = Pt(sw_px * 0.75)
                except Exception:
                    pass
            continue

