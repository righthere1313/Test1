import hashlib
from pathlib import Path
from typing import Any, Dict, Optional

from pptx import Presentation


def _assets_dir_for_template(template_path: Path) -> Path:
    h = hashlib.md5(str(template_path).encode("utf-8")).hexdigest()[:12]
    return Path("data/generated/ppt/assets").resolve() / "template_extract" / h


def _safe_text(text: str, max_len: int) -> str:
    t = (text or "").strip()
    if len(t) <= max_len:
        return t
    return t[: max(0, int(max_len) - 1)] + "…"


def _placeholder_info(shape) -> Optional[dict]:
    try:
        if not getattr(shape, "is_placeholder", False):
            return None
        pf = getattr(shape, "placeholder_format", None)
        idx = getattr(pf, "idx", None)
        ph_type = getattr(getattr(pf, "type", None), "name", None)
        return {"idx": idx, "type": ph_type}
    except Exception:
        return None


def _shape_rect(shape) -> Optional[dict]:
    try:
        return {
            "left": int(getattr(shape, "left", 0)),
            "top": int(getattr(shape, "top", 0)),
            "width": int(getattr(shape, "width", 0)),
            "height": int(getattr(shape, "height", 0)),
        }
    except Exception:
        return None


def _extract_image_asset(shape, out_dir: Path) -> Optional[dict]:
    try:
        image = getattr(shape, "image", None)
        if image is None:
            return None
        blob = getattr(image, "blob", None)
        if not blob:
            return None
        ext = str(getattr(image, "ext", "") or "").lstrip(".").lower() or "png"
        if ext not in {"png", "jpg", "jpeg", "gif", "bmp", "tif", "tiff", "webp"}:
            ext = "png"
        key = hashlib.sha1(blob).hexdigest()[:16]
        out_dir.mkdir(parents=True, exist_ok=True)
        out_path = out_dir / f"img_{key}.{ext}"
        if not out_path.exists():
            out_path.write_bytes(blob)
        return {"image_path": str(out_path), "ext": ext, "size_bytes": int(len(blob)), "sha1": key}
    except Exception:
        return None


def parse_ppt_template(
    template_path: Path,
    *,
    max_slides: int = 12,
    max_shapes_per_slide: int = 80,
    max_text_len: int = 400,
) -> Dict[str, Any]:
    prs = Presentation(str(template_path))
    out_dir = _assets_dir_for_template(template_path)

    def iter_shapes(root_shapes):
        for s in list(root_shapes or []):
            yield s
            try:
                if getattr(s, "shape_type", None) == 6 and hasattr(s, "shapes"):
                    for child in iter_shapes(getattr(s, "shapes", [])):
                        yield child
            except Exception:
                continue

    layouts: list[dict] = []
    for idx, layout in enumerate(prs.slide_layouts):
        item: dict = {"layout_index": idx, "name": getattr(layout, "name", "") or "", "placeholders": []}
        try:
            for ph in list(layout.placeholders):
                rect = _shape_rect(ph) or {"left": 0, "top": 0, "width": 0, "height": 0}
                info = _placeholder_info(ph) or {}
                item["placeholders"].append(
                    {
                        "idx": info.get("idx"),
                        "type": info.get("type"),
                        **rect,
                    }
                )
        except Exception:
            pass
        layouts.append(item)

    slides: list[dict] = []
    media: dict[str, dict] = {}
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    for i, slide in enumerate(prs.slides, start=1):
        if i > int(max_slides):
            break
        slide_item: dict = {"slide_number": i, "layout_name": "", "shapes": []}
        try:
            layout = getattr(slide, "slide_layout", None)
            slide_item["layout_name"] = getattr(layout, "name", "") or ""
        except Exception:
            slide_item["layout_name"] = ""

        collected = 0
        for shp in iter_shapes(getattr(slide, "shapes", [])):
            if collected >= int(max_shapes_per_slide):
                break
            rect = _shape_rect(shp)
            if rect is None:
                continue
            kind = "shape"
            if hasattr(shp, "text_frame"):
                kind = "text"
            if getattr(shp, "shape_type", None) == 13:
                kind = "picture"

            item: dict = {"kind": kind, **rect}
            ph = _placeholder_info(shp)
            if ph:
                item["placeholder"] = ph

            if kind == "text":
                try:
                    item["text"] = _safe_text(getattr(shp.text_frame, "text", "") or "", max_text_len)
                except Exception:
                    item["text"] = ""
            if kind == "picture":
                asset = _extract_image_asset(shp, out_dir)
                if asset:
                    item["image"] = {"image_path": asset["image_path"], "ext": asset["ext"], "size_bytes": asset["size_bytes"]}
                    media.setdefault(asset["sha1"], asset)
                    if rect["left"] == 0 and rect["top"] == 0 and rect["width"] == sw and rect["height"] == sh:
                        item["is_full_slide_background"] = True

            slide_item["shapes"].append(item)
            collected += 1

        slides.append(slide_item)

    media_list = sorted(media.values(), key=lambda x: x["sha1"])
    return {
        "template_path": str(template_path),
        "template_name": template_path.name,
        "slide_width": sw,
        "slide_height": sh,
        "layout_count": len(prs.slide_layouts),
        "layouts": layouts,
        "slides": slides,
        "media": media_list,
        "assets_dir": str(out_dir),
        "max_slides": int(max_slides),
        "max_shapes_per_slide": int(max_shapes_per_slide),
    }
