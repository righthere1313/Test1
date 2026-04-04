"""
PPT 生成与模板适配工具。

包含能力：
- 模板解析：从 /ppt模版 目录选择模板并清空模板自带页面，仅保留主题/母版资源。
- 背景图：由上游（大模型/客户端）在 JSON 中显式指定 background_image。
- 插图：支持 image_path / image_filename（multipart 上传映射）插入图片，并自动排版避免文字覆盖。
- 插图缩放：按页面约 1/4 面积目标等比缩放，在可用图片区内居中放置。
"""

import uuid
import os
import hashlib
from typing import Optional
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from app.schemas.generation import PPTPresentation
from app.services.svg_layouts import read_layout_svg
from app.services.svg_render import extract_svg_token_styles, render_svg_to_png_file, strip_svg_tokens
from app.services.svg_to_ppt import draw_svg_visuals
from app.services.charts_service import render_chart_png


def _apply_background(slide, prs: Presentation, image_path: Optional[str]) -> None:
    if not image_path:
        return
    try:
        p = _resolve_local_path(image_path)
        if p is None:
            return
        pic = slide.shapes.add_picture(str(p), 0, 0, width=prs.slide_width, height=prs.slide_height)
        sp_tree = slide.shapes._spTree
        sp_tree.remove(pic._element)
        sp_tree.insert(2, pic._element)
    except Exception:
        return


def _resolve_local_path(path_like: Optional[str]) -> Optional[Path]:
    raw = str(path_like or "").strip()
    if not raw:
        return None
    if raw.startswith("file://"):
        raw = raw[len("file://") :]
    p = Path(raw)
    if not p.is_absolute():
        repo_root = Path(__file__).resolve().parents[4]
        p = (repo_root / p).resolve()
    if not p.exists() or not p.is_file():
        return None
    return p


def _assets_dir() -> Path:
    return Path("data/generated/ppt/assets").resolve()


def _hex_color(s: str) -> Optional[RGBColor]:
    raw = (s or "").strip()
    if not raw:
        return None
    if raw.startswith("#"):
        raw = raw[1:]
    if len(raw) == 3:
        raw = "".join([c * 2 for c in raw])
    if len(raw) != 6:
        return None
    try:
        r = int(raw[0:2], 16)
        g = int(raw[2:4], 16)
        b = int(raw[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return None


def _pick_font_name(family: str) -> str:
    s = (family or "").strip()
    if not s:
        return "Microsoft YaHei"
    first = s.split(",", 1)[0].strip()
    first = first.strip('"').strip("'").strip()
    return first or "Microsoft YaHei"


def _align_from_anchor(anchor: str):
    a = (anchor or "").strip().lower()
    if a in {"middle", "center"}:
        return PP_ALIGN.CENTER
    if a in {"end", "right"}:
        return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


def _add_styled_textbox(
    slide,
    *,
    left: int,
    top: int,
    width: int,
    height: int,
    lines: list[str],
    font_name: str,
    font_size_pt: float,
    color_hex: str,
    bold: bool,
    align,
) -> Optional[object]:
    if width <= 0 or height <= 0:
        return None
    try:
        shp = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    except Exception:
        return None
    try:
        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True
    except Exception:
        return shp
    if not lines:
        return shp
    tf.text = str(lines[0] or "")
    for x in lines[1:]:
        p = tf.add_paragraph()
        p.text = str(x or "")
        p.level = 0
    rgb = _hex_color(color_hex)
    for p in getattr(tf, "paragraphs", []):
        try:
            p.alignment = align
        except Exception:
            pass
        try:
            p.space_before = Pt(0)
            p.space_after = Pt(0)
        except Exception:
            pass
        for r in getattr(p, "runs", []):
            try:
                r.font.name = font_name
            except Exception:
                pass
            try:
                r.font.size = Pt(float(font_size_pt))
            except Exception:
                pass
            try:
                r.font.bold = bool(bold)
            except Exception:
                pass
            if rgb is not None:
                try:
                    r.font.color.rgb = rgb
                except Exception:
                    pass
    return shp


def _split_lines_to_columns(lines: list[str], box_h: int, box_w: int, *, font_pt: float) -> list[list[str]]:
    pts = float(font_pt)
    if pts <= 0:
        pts = 24.0
    line_h = int(Pt(pts * 1.25))
    max_lines = max(1, int(max(1, box_h) // max(1, line_h)))
    if not lines:
        return []
    need = int((len(lines) + max_lines - 1) // max_lines)
    max_cols = 4
    try:
        fit = int(max(1, box_w) // int(Inches(2.8)))
        fit = max(1, min(max_cols, fit))
    except Exception:
        fit = max_cols
    cols = max(1, min(need, fit))
    per_col = max_lines
    out: list[list[str]] = [[] for _ in range(cols)]
    idx = 0
    for c in range(cols):
        if idx >= len(lines):
            break
        out[c] = lines[idx : idx + per_col]
        idx += per_col
    while idx < len(lines):
        out[-1].append(lines[idx])
        idx += 1
    return out


def _split_rect_for_image(rect: tuple[int, int, int, int], position: str) -> tuple[tuple[int, int, int, int], Optional[tuple[int, int, int, int]]]:
    left, top, width, height = rect
    pos = (position or "").strip().lower()
    if pos == "full":
        return (left, top, width, height), (left, top, width, height)
    if pos == "left":
        img_w = int(width * 0.42)
        gap = int(Inches(0.25))
        img_rect = (left, top, img_w, height)
        text_rect = (left + img_w + gap, top, max(1, width - img_w - gap), height)
        return text_rect, img_rect
    if pos == "right":
        img_w = int(width * 0.42)
        gap = int(Inches(0.25))
        img_rect = (left + width - img_w, top, img_w, height)
        text_rect = (left, top, max(1, width - img_w - gap), height)
        return text_rect, img_rect
    if pos == "top":
        img_h = int(height * 0.55)
        gap = int(Inches(0.20))
        img_rect = (left, top, width, img_h)
        text_rect = (left, top + img_h + gap, width, max(1, height - img_h - gap))
        return text_rect, img_rect
    if pos == "bottom":
        img_h = int(height * 0.55)
        gap = int(Inches(0.20))
        img_rect = (left, top + height - img_h, width, img_h)
        text_rect = (left, top, width, max(1, height - img_h - gap))
        return text_rect, img_rect
    return (left, top, width, height), None


def _normalize_part_rect(part: dict, slide_w: int, slide_h: int) -> Optional[tuple[int, int, int, int]]:
    try:
        unit = str(part.get("unit") or "emu").strip().lower()
        x = float(part.get("x"))
        y = float(part.get("y"))
        w = float(part.get("w"))
        h = float(part.get("h"))
        if unit == "ratio":
            return int(x * slide_w), int(y * slide_h), int(w * slide_w), int(h * slide_h)
        return int(x), int(y), int(w), int(h)
    except Exception:
        return None


def _normalize_rect_spec(spec: Optional[object], slide_w: int, slide_h: int) -> Optional[tuple[int, int, int, int]]:
    if spec is None:
        return None
    try:
        if hasattr(spec, "model_dump"):
            d = spec.model_dump()
        elif isinstance(spec, dict):
            d = spec
        else:
            return None
    except Exception:
        return None
    try:
        unit = str(d.get("unit") or "emu").strip().lower()
        x = float(d.get("x"))
        y = float(d.get("y"))
        w = float(d.get("w"))
        h = float(d.get("h"))
        if unit == "ratio":
            return int(x * slide_w), int(y * slide_h), int(w * slide_w), int(h * slide_h)
        return int(x), int(y), int(w), int(h)
    except Exception:
        return None


def _compose_background_from_parts(prs: Presentation, parts: list[dict]) -> Optional[str]:
    if not parts:
        return None
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    normalized = []
    for p in parts:
        if not isinstance(p, dict):
            continue
        img = _resolve_local_path(p.get("image_path"))
        if img is None:
            continue
        rect = _normalize_part_rect(p, slide_w, slide_h)
        if rect is None:
            continue
        z = int(p.get("z") or 0)
        normalized.append((z, img, rect))
    if not normalized:
        return None

    try:
        from PIL import Image
    except Exception:
        return None

    normalized.sort(key=lambda x: x[0])
    scale = 1920.0 / float(max(1, slide_w))
    out_w = 1920
    out_h = max(1, int(round(float(slide_h) * scale)))
    canvas = Image.new("RGBA", (out_w, out_h), (255, 255, 255, 0))

    for _, img_path, (x, y, w, h) in normalized:
        try:
            with Image.open(str(img_path)) as im:
                im_rgba = im.convert("RGBA")
                x_px = int(round(float(x) * scale))
                y_px = int(round(float(y) * scale))
                w_px = max(1, int(round(float(w) * scale)))
                h_px = max(1, int(round(float(h) * scale)))
                im_resized = im_rgba.resize((w_px, h_px))
                canvas.alpha_composite(im_resized, (x_px, y_px))
        except Exception:
            continue

    try:
        out_dir = _assets_dir()
        out_dir.mkdir(parents=True, exist_ok=True)
        key_src = "|".join([f"{z}:{str(p)}:{r[0]}:{r[1]}:{r[2]}:{r[3]}" for z, p, r in normalized]).encode("utf-8")
        key = hashlib.md5(key_src).hexdigest()[:16]
        out_path = out_dir / f"bg_composed_{key}.png"
        if not out_path.exists():
            canvas.save(out_path)
        return str(out_path)
    except Exception:
        return None


def _apply_background_parts(slide, prs: Presentation, parts: list[dict]) -> None:
    if not parts:
        return
    composed = _compose_background_from_parts(prs, parts)
    if composed:
        _apply_background(slide, prs, composed)
        return
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    items = []
    for p in parts:
        if not isinstance(p, dict):
            continue
        img = _resolve_local_path(p.get("image_path"))
        if img is None:
            continue
        rect = _normalize_part_rect(p, slide_w, slide_h)
        if rect is None:
            continue
        z = int(p.get("z") or 0)
        items.append((z, img, rect))
    if not items:
        return
    items.sort(key=lambda x: x[0])
    try:
        sp_tree = slide.shapes._spTree
    except Exception:
        sp_tree = None
    for _, img, (x, y, w, h) in items:
        try:
            pic = slide.shapes.add_picture(str(img), x, y, width=w, height=h)
            if sp_tree is not None:
                try:
                    sp_tree.remove(pic._element)
                    sp_tree.insert(2, pic._element)
                except Exception:
                    pass
        except Exception:
            continue


def _shape_spid(shape) -> Optional[str]:
    try:
        c_nv_pr = shape._element.xpath('.//*[local-name()="cNvPr"]')
        if not c_nv_pr:
            return None
        spid = c_nv_pr[0].get("id")
        return str(spid) if spid else None
    except Exception:
        return None


def _set_font_size_for_text_frame(text_frame, size_pt: int) -> None:
    try:
        for p in text_frame.paragraphs:
            for r in getattr(p, "runs", []):
                try:
                    r.font.size = Pt(size_pt)
                except Exception:
                    continue
    except Exception:
        return


def _get_effective_font_size_pt(text_frame, default_pt: int = 18) -> int:
    try:
        for p in text_frame.paragraphs:
            for r in getattr(p, "runs", []):
                sz = getattr(r.font, "size", None)
                if sz is not None:
                    return int(sz.pt)
    except Exception:
        pass
    return int(default_pt)


def _should_shrink_text(points: list[str], box_width_emu: int, box_height_emu: int) -> bool:
    try:
        if box_width_emu <= 0 or box_height_emu <= 0:
            return False
        total_chars = sum(len(str(t or "")) for t in points)
        bullets = max(1, len(points))
        score = total_chars + bullets * 20
        w_in = float(box_width_emu) / 914400.0
        h_in = float(box_height_emu) / 914400.0
        capacity = int(max(80.0, (w_in * h_in) * 20.0))
        return score > capacity
    except Exception:
        return False


def _rect_overlap(a: tuple[int, int, int, int], b: tuple[int, int, int, int]) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def _shape_rect(shape) -> Optional[tuple[int, int, int, int]]:
    try:
        return int(shape.left), int(shape.top), int(shape.width), int(shape.height)
    except Exception:
        return None


def _remove_empty_overlapping_text_shapes(slide, keep_shapes: list) -> None:
    keep_ids = {id(s) for s in keep_shapes if s is not None}
    keep_rects = [r for s in keep_shapes if s is not None for r in [_shape_rect(s)] if r is not None]
    if not keep_rects:
        return
    try:
        sp_tree = slide.shapes._spTree
    except Exception:
        return
    to_remove = []
    for shp in getattr(slide, "shapes", []):
        if id(shp) in keep_ids:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        try:
            if (shp.text_frame.text or "").strip():
                continue
        except Exception:
            continue
        r = _shape_rect(shp)
        if r is None:
            continue
        if any(_rect_overlap(r, kr) for kr in keep_rects):
            try:
                to_remove.append(shp._element)
            except Exception:
                continue
    for el in to_remove:
        try:
            sp_tree.remove(el)
        except Exception:
            continue


def _add_single_click_reveal(slide, shapes: list) -> None:
    spids = []
    for s in shapes:
        spid = _shape_spid(s)
        if spid:
            spids.append(spid)
    if not spids:
        return
    try:
        for child in list(slide._element):
            if child.tag.endswith("}timing"):
                return
    except Exception:
        pass
    try:
        timing = OxmlElement("p:timing")

        tn_lst = OxmlElement("p:tnLst")
        timing.append(tn_lst)

        par_root = OxmlElement("p:par")
        tn_lst.append(par_root)

        ctn_root = OxmlElement("p:cTn")
        ctn_root.set("id", "1")
        ctn_root.set("dur", "indefinite")
        ctn_root.set("restart", "never")
        ctn_root.set("nodeType", "tmRoot")
        par_root.append(ctn_root)

        root_child = OxmlElement("p:childTnLst")
        ctn_root.append(root_child)

        seq = OxmlElement("p:seq")
        seq.set("concurrent", "1")
        seq.set("nextAc", "seek")
        root_child.append(seq)

        ctn_seq = OxmlElement("p:cTn")
        ctn_seq.set("id", "2")
        ctn_seq.set("dur", "indefinite")
        ctn_seq.set("nodeType", "mainSeq")
        seq.append(ctn_seq)

        seq_child = OxmlElement("p:childTnLst")
        ctn_seq.append(seq_child)

        par_effect = OxmlElement("p:par")
        seq_child.append(par_effect)

        ctn_effect = OxmlElement("p:cTn")
        ctn_effect.set("id", "3")
        ctn_effect.set("presetID", "10")
        ctn_effect.set("presetClass", "entr")
        ctn_effect.set("presetSubtype", "0")
        ctn_effect.set("fill", "hold")
        ctn_effect.set("grpId", "0")
        ctn_effect.set("nodeType", "clickEffect")
        par_effect.append(ctn_effect)

        st_cond = OxmlElement("p:stCondLst")
        ctn_effect.append(st_cond)
        cond = OxmlElement("p:cond")
        cond.set("delay", "0")
        st_cond.append(cond)

        effect_child = OxmlElement("p:childTnLst")
        ctn_effect.append(effect_child)

        next_id = 10
        for spid in spids:
            set_el = OxmlElement("p:set")
            effect_child.append(set_el)
            bhvr = OxmlElement("p:cBhvr")
            set_el.append(bhvr)
            ctn_set = OxmlElement("p:cTn")
            ctn_set.set("id", str(next_id))
            next_id += 1
            ctn_set.set("dur", "1")
            ctn_set.set("fill", "hold")
            bhvr.append(ctn_set)
            st2 = OxmlElement("p:stCondLst")
            ctn_set.append(st2)
            cond2 = OxmlElement("p:cond")
            cond2.set("delay", "0")
            st2.append(cond2)
            tgt = OxmlElement("p:tgtEl")
            bhvr.append(tgt)
            sp_tgt = OxmlElement("p:spTgt")
            sp_tgt.set("spid", str(spid))
            tgt.append(sp_tgt)
            attr = OxmlElement("p:attrNameLst")
            bhvr.append(attr)
            attr_name = OxmlElement("p:attrName")
            attr_name.text = "style.visibility"
            attr.append(attr_name)
            to = OxmlElement("p:to")
            set_el.append(to)
            sval = OxmlElement("p:strVal")
            sval.set("val", "visible")
            to.append(sval)

            anim = OxmlElement("p:animEffect")
            anim.set("transition", "in")
            anim.set("filter", "fade")
            effect_child.append(anim)
            bhvr2 = OxmlElement("p:cBhvr")
            anim.append(bhvr2)
            ctn_anim = OxmlElement("p:cTn")
            ctn_anim.set("id", str(next_id))
            next_id += 1
            ctn_anim.set("dur", "500")
            bhvr2.append(ctn_anim)
            tgt2 = OxmlElement("p:tgtEl")
            bhvr2.append(tgt2)
            sp_tgt2 = OxmlElement("p:spTgt")
            sp_tgt2.set("spid", str(spid))
            tgt2.append(sp_tgt2)

        prev_lst = OxmlElement("p:prevCondLst")
        seq.append(prev_lst)
        prev = OxmlElement("p:cond")
        prev.set("evt", "onPrev")
        prev.set("delay", "0")
        prev_lst.append(prev)
        prev_tgt = OxmlElement("p:tgtEl")
        prev.append(prev_tgt)
        prev_tgt.append(OxmlElement("p:sldTgt"))

        next_lst = OxmlElement("p:nextCondLst")
        seq.append(next_lst)
        nxt = OxmlElement("p:cond")
        nxt.set("evt", "onNext")
        nxt.set("delay", "0")
        next_lst.append(nxt)
        nxt_tgt = OxmlElement("p:tgtEl")
        nxt.append(nxt_tgt)
        nxt_tgt.append(OxmlElement("p:sldTgt"))

        bld_lst = OxmlElement("p:bldLst")
        timing.append(bld_lst)
        for spid in spids:
            bld = OxmlElement("p:bldP")
            bld.set("spid", str(spid))
            bld.set("grpId", "0")
            bld_lst.append(bld)

        slide._element.append(timing)
    except Exception:
        return


def _insert_picture(slide, prs: Presentation, image_path: Optional[str], position: str) -> None:
    if not image_path:
        return
    try:
        pos = (position or "").strip().lower()
        text_rect, img_rect = _layout_rects_for_image(slide, prs, pos)
        if img_rect is None:
            img_rect = text_rect
        _place_picture_in_rect(slide, prs, image_path, img_rect, pos)
    except Exception:
        return


def _image_aspect_ratio(image_path: str) -> float:
    try:
        from PIL import Image

        with Image.open(image_path) as im:
            w, h = im.size
        if w <= 0 or h <= 0:
            return 4.0 / 3.0
        return float(w) / float(h)
    except Exception:
        return 4.0 / 3.0


def _compute_size_for_target_area(aspect: float, target_area: float, max_w: int, max_h: int) -> tuple[int, int]:
    if aspect <= 0:
        aspect = 4.0 / 3.0
    w = (target_area * aspect) ** 0.5
    h = (target_area / aspect) ** 0.5
    if w <= 0 or h <= 0:
        return max(1, max_w), max(1, max_h)
    scale = min(float(max_w) / w, float(max_h) / h, 1e9)
    if scale < 1.0:
        w *= scale
        h *= scale
    if w > max_w:
        w = max_w
    if h > max_h:
        h = max_h
    return max(1, int(w)), max(1, int(h))


def _place_picture_in_rect(slide, prs: Presentation, image_path: str, rect: tuple[int, int, int, int], position: str) -> None:
    l, t, w, h = rect
    if (position or "").strip().lower() == "full":
        slide.shapes.add_picture(image_path, l, t, width=w, height=h)
        return

    padding = int(Inches(0.25))
    max_w = max(1, w - padding * 2)
    max_h = max(1, h - padding * 2)
    target_area = float(prs.slide_width) * float(prs.slide_height) * 0.25
    aspect = _image_aspect_ratio(image_path)
    pic_w, pic_h = _compute_size_for_target_area(aspect, target_area, max_w, max_h)
    left = l + int((w - pic_w) / 2)
    top = t + int((h - pic_h) / 2)
    slide.shapes.add_picture(image_path, left, top, width=pic_w, height=pic_h)


def _body_rect(slide, prs: Presentation) -> tuple[int, int, int, int]:
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.6)
    default_top = Inches(1.7)
    gap = int(Inches(0.2))
    header_bottom = 0
    footer_top: Optional[int] = None
    try:
        header_zone = int(sh * 0.42)
        footer_zone = int(sh * 0.78)

        title_shape = None
        try:
            title_shape = getattr(slide.shapes, "title", None)
        except Exception:
            title_shape = None
        if title_shape is not None:
            try:
                header_bottom = max(header_bottom, int(title_shape.top) + int(title_shape.height))
            except Exception:
                pass

        for shp in slide.shapes:
            try:
                top = int(getattr(shp, "top", 0))
                height = int(getattr(shp, "height", 0))
                bottom = top + height
            except Exception:
                continue

            if top <= header_zone and height <= int(sh * 0.45):
                if hasattr(shp, "text_frame") or getattr(shp, "is_placeholder", False):
                    header_bottom = max(header_bottom, bottom)

            if top >= footer_zone and height <= int(sh * 0.25):
                if hasattr(shp, "text_frame") or getattr(shp, "is_placeholder", False):
                    if footer_top is None or top < footer_top:
                        footer_top = top
    except Exception:
        header_bottom = 0
        footer_top = None
    body_top = max(int(default_top), header_bottom + gap)
    left = int(margin)
    top = int(body_top)
    width = int(sw - margin * 2)
    max_bottom = int(sh - margin)
    if footer_top is not None:
        max_bottom = min(max_bottom, max(int(body_top + gap), int(footer_top - gap)))
    height = max(1, int(max_bottom - body_top))
    return left, top, width, height


def _layout_rects_for_image(slide, prs: Presentation, position: str) -> tuple[tuple[int, int, int, int], Optional[tuple[int, int, int, int]]]:
    left, top, width, height = _body_rect(slide, prs)
    pos = (position or "").strip().lower()
    if pos == "full":
        return (left, top, width, height), (0, 0, int(prs.slide_width), int(prs.slide_height))
    if pos == "left":
        img_w = int(width * 0.55)
        gap = int(Inches(0.35))
        img_rect = (left, top, img_w, height)
        text_rect = (left + img_w + gap, top, width - img_w - gap, height)
        return text_rect, img_rect
    if pos == "right":
        img_w = int(width * 0.55)
        gap = int(Inches(0.35))
        img_rect = (left + width - img_w, top, img_w, height)
        text_rect = (left, top, width - img_w - gap, height)
        return text_rect, img_rect
    if pos == "top":
        img_h = int(height * 0.90)
        gap = int(Inches(0.20))
        img_rect = (left, top, width, img_h)
        text_rect = (left, top + img_h + gap, width, height - img_h - gap)
        return text_rect, img_rect
    if pos == "bottom":
        img_h = int(height * 0.90)
        gap = int(Inches(0.20))
        img_rect = (left, top + height - img_h, width, img_h)
        text_rect = (left, top, width, height - img_h - gap)
        return text_rect, img_rect
    return (left, top, width, height), None


def _ensure_textbox(slide, rect: tuple[int, int, int, int]):
    l, t, w, h = rect
    try:
        return slide.shapes.add_textbox(l, t, w, h)
    except Exception:
        return None


def _move_shape_to_rect(shape, rect: tuple[int, int, int, int]) -> None:
    l, t, w, h = rect
    try:
        shape.left = l
        shape.top = t
        shape.width = w
        shape.height = h
    except Exception:
        return


def _set_subtitle_if_possible(slide, text: Optional[str]):
    if not text:
        return None
    try:
        slide.placeholders[1].text = text
        return slide.placeholders[1]
    except Exception:
        pass
    for shape in getattr(slide, "placeholders", []):
        try:
            if hasattr(shape, "text"):
                if getattr(shape, "text", "") == "":
                    shape.text = text
                    return shape
        except Exception:
            continue
    try:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(0.7))
        tb.text_frame.text = text
        return tb
    except Exception:
        return None


def _set_title_if_possible(slide, text: str):
    if not text:
        return None
    try:
        title_shape = getattr(slide.shapes, "title", None)
        if title_shape is not None:
            title_shape.text = text
            return title_shape
    except Exception:
        pass
    for shape in getattr(slide, "placeholders", []):
        try:
            if hasattr(shape, "text"):
                shape.text = text
                return shape
            if hasattr(shape, "text_frame"):
                shape.text_frame.text = text
                return shape
        except Exception:
            continue
    for shape in getattr(slide, "shapes", []):
        try:
            if hasattr(shape, "text"):
                shape.text = text
                return shape
            if hasattr(shape, "text_frame"):
                shape.text_frame.text = text
                return shape
        except Exception:
            continue
    try:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(8.4), Inches(0.9))
        tb.text_frame.text = text
        return tb
    except Exception:
        return None


def _get_body_shape(slide):
    for shape in getattr(slide, "placeholders", []):
        if getattr(shape, "is_placeholder", False) and hasattr(shape, "text_frame"):
            if getattr(getattr(shape, "placeholder_format", None), "idx", None) == 0:
                continue
            return shape
    for shape in getattr(slide, "shapes", []):
        if hasattr(shape, "text_frame") and shape is not getattr(slide.shapes, "title", None):
            return shape
    return None


def _get_best_body_shape(slide):
    title_shape = None
    try:
        title_shape = getattr(slide.shapes, "title", None)
    except Exception:
        title_shape = None

    best = None
    best_area = -1
    preferred_types = {"BODY", "OBJECT", "CONTENT", "TEXT"}
    forbidden_types = {"FOOTER", "DATE", "SLIDE_NUMBER", "HEADER"}

    for shape in getattr(slide, "placeholders", []):
        if not getattr(shape, "is_placeholder", False) or not hasattr(shape, "text_frame"):
            continue
        pf = getattr(shape, "placeholder_format", None)
        try:
            if getattr(pf, "idx", None) == 0:
                continue
        except Exception:
            continue
        try:
            ph_type = str(getattr(getattr(pf, "type", None), "name", "") or "").upper()
        except Exception:
            ph_type = ""
        if ph_type in forbidden_types:
            continue
        try:
            area = int(shape.width) * int(shape.height)
        except Exception:
            area = 0
        boost = 1_000_000_000 if ph_type in preferred_types else 0
        score = boost + area
        if score > best_area:
            best = shape
            best_area = score
    if best is not None:
        return best

    for shape in getattr(slide, "shapes", []):
        if shape is title_shape:
            continue
        if not hasattr(shape, "text_frame"):
            continue
        try:
            area = int(shape.width) * int(shape.height)
        except Exception:
            area = 0
        if area > best_area:
            best = shape
            best_area = area
    return best


def _find_picture_placeholder(slide):
    for shape in getattr(slide, "placeholders", []):
        if not getattr(shape, "is_placeholder", False):
            continue
        pf = getattr(shape, "placeholder_format", None)
        try:
            ph_type = getattr(getattr(pf, "type", None), "name", "") or ""
        except Exception:
            ph_type = ""
        if ph_type.upper() == "PICTURE":
            return shape
    return None


def _fill_text_frame(tf, lines: list[str]) -> None:
    try:
        tf.clear()
    except Exception:
        return
    try:
        tf.word_wrap = True
    except Exception:
        pass
    if not lines:
        return
    tf.text = str(lines[0] or "")
    for point in lines[1:]:
        p = tf.add_paragraph()
        p.text = str(point or "")
        p.level = 0


def _el_id(shape) -> Optional[int]:
    if shape is None:
        return None
    try:
        return id(shape._element)
    except Exception:
        return None


def _iter_all_shapes(container):
    try:
        shapes = list(getattr(container, "shapes", []))
    except Exception:
        shapes = []
    for shp in shapes:
        yield shp
        try:
            if getattr(shp, "shape_type", None) == 6 and hasattr(shp, "shapes"):
                for child in _iter_all_shapes(shp):
                    yield child
        except Exception:
            continue


def _clear_unused_placeholder_text(slide, keep_shapes: list) -> None:
    keep_ids = {x for x in (_el_id(s) for s in keep_shapes) if x is not None}
    for shp in getattr(slide, "placeholders", []):
        if _el_id(shp) in keep_ids:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        try:
            if (shp.text_frame.text or "").strip() == "":
                continue
        except Exception:
            continue
        try:
            shp.text_frame.clear()
        except Exception:
            try:
                shp.text_frame.text = ""
            except Exception:
                continue


def _clear_unkept_text(slide, keep_shapes: list) -> None:
    keep_ids = {x for x in (_el_id(s) for s in keep_shapes) if x is not None}
    for shp in _iter_all_shapes(slide):
        if _el_id(shp) in keep_ids:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        try:
            if (shp.text_frame.text or "").strip() == "":
                continue
        except Exception:
            continue
        try:
            shp.text_frame.clear()
        except Exception:
            try:
                shp.text_frame.text = ""
            except Exception:
                continue


def _bring_to_front(slide, shape) -> None:
    if shape is None:
        return
    try:
        sp_tree = slide.shapes._spTree
        el = shape._element
        sp_tree.remove(el)
        sp_tree.append(el)
    except Exception:
        return


def _shape_area(shape) -> int:
    try:
        return int(shape.width) * int(shape.height)
    except Exception:
        return 0


def _candidate_text_shapes(slide, prs: Presentation, exclude: list) -> list:
    exclude_ids = {x for x in (_el_id(s) for s in exclude) if x is not None}
    slide_area = int(prs.slide_width) * int(prs.slide_height)
    min_area = max(1, int(slide_area * 0.001))
    out = []
    for shp in _iter_all_shapes(slide):
        if _el_id(shp) in exclude_ids:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        if _shape_area(shp) < min_area:
            continue
        out.append(shp)
    out.sort(key=_shape_area, reverse=True)
    return out


def _fill_text_shapes(shapes: list, lines: list[str]) -> list:
    filled = []
    shapes = [s for s in shapes if s is not None and hasattr(s, "text_frame")]
    if not shapes:
        return filled
    lines = [str(x or "").strip() for x in (lines or [])]
    lines = [x for x in lines if x]
    if not lines:
        return filled

    max_boxes = min(3, len(shapes))
    targets = shapes[:max_boxes]
    per = max(1, int((len(lines) + max_boxes - 1) / max_boxes))
    idx = 0
    for s in targets:
        chunk = lines[idx : idx + per]
        idx += per
        if not chunk:
            continue
        try:
            _fill_text_frame(s.text_frame, chunk)
            try:
                if _should_shrink_text(chunk, int(s.width), int(s.height)):
                    base = _get_effective_font_size_pt(s.text_frame, default_pt=18)
                    _set_font_size_for_text_frame(s.text_frame, max(10, base - 2))
            except Exception:
                pass
            filled.append(s)
        except Exception:
            continue
    return filled


def _fill_text_blocks(shapes: list, blocks: list[list[str]]) -> list:
    filled = []
    shapes = [s for s in shapes if s is not None and hasattr(s, "text_frame")]
    if not shapes:
        return filled
    idx = 0
    for block in blocks:
        if idx >= len(shapes):
            break
        if not isinstance(block, list):
            continue
        s = shapes[idx]
        idx += 1
        lines = [str(x or "").strip() for x in block]
        lines = [x for x in lines if x]
        if not lines:
            try:
                s.text_frame.text = ""
            except Exception:
                pass
            continue
        try:
            _fill_text_frame(s.text_frame, lines)
            try:
                if _should_shrink_text(lines, int(s.width), int(s.height)):
                    base = _get_effective_font_size_pt(s.text_frame, default_pt=18)
                    _set_font_size_for_text_frame(s.text_frame, max(10, base - 2))
            except Exception:
                pass
            filled.append(s)
        except Exception:
            continue
    return filled


class PPTGenerator:
    OUTPUT_DIR = "data/generated/ppt"

    def __init__(self) -> None:
        os.makedirs(self.OUTPUT_DIR, exist_ok=True)

    def generate(self, data: PPTPresentation, image_files: Optional[dict[str, str]] = None) -> str:
        layout_name = str(getattr(data, "layout", None) or "").strip()
        if layout_name:
            prs = Presentation()
            try:
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
            except Exception:
                pass
            slide_w = int(prs.slide_width)
            slide_h = int(prs.slide_height)
            rect = _normalize_rect_spec(getattr(data, "content_rect_default", None), slide_w, slide_h)

            try:
                blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[len(prs.slide_layouts) - 1]
            except Exception:
                blank_layout = prs.slide_layouts[0]

            cover_svg = read_layout_svg(layout_name, "01_cover.svg")
            content_svg = read_layout_svg(layout_name, "03_content.svg")
            ending_svg = read_layout_svg(layout_name, "04_ending.svg")

            cover_styles = extract_svg_token_styles(cover_svg)
            content_styles = extract_svg_token_styles(content_svg)
            ending_styles = extract_svg_token_styles(ending_svg)

            out_dir = _assets_dir()
            cover_bg = render_svg_to_png_file(strip_svg_tokens(cover_svg), out_dir=out_dir, out_prefix=f"layout_{layout_name}_cover", width_px=1920)
            content_bg = render_svg_to_png_file(strip_svg_tokens(content_svg), out_dir=out_dir, out_prefix=f"layout_{layout_name}_content", width_px=1920)
            ending_bg = render_svg_to_png_file(strip_svg_tokens(ending_svg), out_dir=out_dir, out_prefix=f"layout_{layout_name}_ending", width_px=1920)

            cover = prs.slides.add_slide(blank_layout)
            if cover_bg:
                _apply_background(cover, prs, cover_bg)
            else:
                draw_svg_visuals(cover, prs, strip_svg_tokens(cover_svg))
            t = str(getattr(data, "title", "") or "").strip()
            st = str(getattr(data, "subtitle", "") or "").strip()
            if "TITLE" in cover_styles and t:
                style, (vw, vh) = cover_styles["TITLE"]
                x = int(style.x / vw * slide_w)
                y = int(style.y / vh * slide_h)
                font_pt = float(style.font_size_px) * 0.75
                width = int(slide_w * 0.86)
                height = int(Pt(font_pt * 1.6))
                left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                top = max(0, int(y - int(Pt(font_pt * 1.1))))
                _add_styled_textbox(
                    cover,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    lines=[t],
                    font_name=_pick_font_name(style.font_family),
                    font_size_pt=font_pt,
                    color_hex=style.fill,
                    bold=style.bold,
                    align=_align_from_anchor(style.anchor),
                )
            if "SUBTITLE" in cover_styles and st:
                style, (vw, vh) = cover_styles["SUBTITLE"]
                x = int(style.x / vw * slide_w)
                y = int(style.y / vh * slide_h)
                font_pt = float(style.font_size_px) * 0.75
                width = int(slide_w * 0.86)
                height = int(Pt(font_pt * 1.4))
                left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                top = max(0, int(y - int(Pt(font_pt * 1.0))))
                _add_styled_textbox(
                    cover,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    lines=[st],
                    font_name=_pick_font_name(style.font_family),
                    font_size_pt=font_pt,
                    color_hex=style.fill,
                    bold=style.bold,
                    align=_align_from_anchor(style.anchor),
                )

            section_name = str(getattr(data, "title", "") or "").strip()
            for i, slide_data in enumerate(getattr(data, "slides", []) or []):
                slide = prs.slides.add_slide(blank_layout)
                if content_bg:
                    _apply_background(slide, prs, content_bg)
                else:
                    draw_svg_visuals(slide, prs, strip_svg_tokens(content_svg))
                page_title = str(getattr(slide_data, "title", "") or "").strip()
                if "PAGE_TITLE" in content_styles and page_title:
                    style, (vw, vh) = content_styles["PAGE_TITLE"]
                    x = int(style.x / vw * slide_w)
                    y = int(style.y / vh * slide_h)
                    font_pt = float(style.font_size_px) * 0.75
                    width = int(slide_w - x - int(slide_w * 0.06))
                    height = int(Pt(font_pt * 1.4))
                    left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                    top = max(0, int(y - int(Pt(font_pt * 1.0))))
                    _add_styled_textbox(
                        slide,
                        left=left,
                        top=top,
                        width=width,
                        height=height,
                        lines=[page_title],
                        font_name=_pick_font_name(style.font_family),
                        font_size_pt=font_pt,
                        color_hex=style.fill,
                        bold=style.bold,
                        align=_align_from_anchor(style.anchor),
                    )
                if rect is not None:
                    base_rect = rect
                    visual_path = None
                    visual_pos = None
                    chart = getattr(slide_data, "chart", None)
                    if chart is not None:
                        try:
                            cd = chart.model_dump() if hasattr(chart, "model_dump") else dict(chart)
                            visual_path = render_chart_png(cd, width_px=1920)
                            visual_pos = str((cd.get("position") or "").strip() or "right")
                        except Exception:
                            visual_path = None
                            visual_pos = None
                    if not visual_path:
                        img_path = getattr(slide_data, "image_path", None)
                        if (not img_path) and getattr(slide_data, "image_filename", None) and image_files:
                            img_path = image_files.get(str(getattr(slide_data, "image_filename") or ""))
                        if img_path:
                            visual_path = str(img_path)
                            visual_pos = str(getattr(slide_data, "image_position", None) or getattr(data, "image_position_default", None) or "right")

                    text_rect, img_rect = _split_rect_for_image(base_rect, visual_pos) if visual_path else (base_rect, None)

                    body_left, body_top, body_w, body_h = text_rect
                    body_lines = [str(x or "").strip() for x in (getattr(slide_data, "content", None) or [])]
                    body_lines = [x for x in body_lines if x]
                    body_font_name = "Microsoft YaHei"
                    body_font_pt = 24.0
                    if "CONTENT_AREA" in content_styles:
                        s0, _ = content_styles["CONTENT_AREA"]
                        body_font_name = _pick_font_name(s0.font_family)
                    cols = _split_lines_to_columns(body_lines, body_h, body_w, font_pt=body_font_pt)
                    if not cols:
                        cols = [[]]
                    gap = int(Inches(0.18)) if len(cols) > 1 else 0
                    total_gap = gap * max(0, len(cols) - 1)
                    col_w = max(1, int((body_w - total_gap) // max(1, len(cols))))
                    for ci, col_lines in enumerate(cols):
                        _add_styled_textbox(
                            slide,
                            left=int(body_left + ci * (col_w + gap)),
                            top=body_top,
                            width=col_w,
                            height=body_h,
                            lines=col_lines,
                            font_name=body_font_name,
                            font_size_pt=body_font_pt,
                            color_hex="#333333",
                            bold=False,
                            align=PP_ALIGN.LEFT,
                        )

                    if visual_path and img_rect is not None:
                        try:
                            _place_picture_in_rect(slide, prs, str(visual_path), img_rect, str(visual_pos or "right"))
                        except Exception:
                            try:
                                _insert_picture(slide, prs, str(visual_path), str(visual_pos or "right"))
                            except Exception:
                                pass
                if "SECTION_NAME" in content_styles and section_name:
                    style, (vw, vh) = content_styles["SECTION_NAME"]
                    x = int(style.x / vw * slide_w)
                    y = int(style.y / vh * slide_h)
                    font_pt = float(style.font_size_px) * 0.75
                    width = int(slide_w * 0.4)
                    height = int(Pt(font_pt * 1.3))
                    left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                    top = max(0, int(y - int(Pt(font_pt * 1.0))))
                    _add_styled_textbox(
                        slide,
                        left=left,
                        top=top,
                        width=width,
                        height=height,
                        lines=[section_name],
                        font_name=_pick_font_name(style.font_family),
                        font_size_pt=font_pt,
                        color_hex=style.fill,
                        bold=style.bold,
                        align=_align_from_anchor(style.anchor),
                    )
                if "PAGE_NUM" in content_styles:
                    style, (vw, vh) = content_styles["PAGE_NUM"]
                    x = int(style.x / vw * slide_w)
                    y = int(style.y / vh * slide_h)
                    font_pt = float(style.font_size_px) * 0.75
                    width = int(slide_w * 0.12)
                    height = int(Pt(font_pt * 1.3))
                    left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                    top = max(0, int(y - int(Pt(font_pt * 1.0))))
                    _add_styled_textbox(
                        slide,
                        left=left,
                        top=top,
                        width=width,
                        height=height,
                        lines=[str(i + 2)],
                        font_name=_pick_font_name(style.font_family),
                        font_size_pt=font_pt,
                        color_hex=style.fill,
                        bold=style.bold,
                        align=_align_from_anchor(style.anchor),
                    )

            ending = prs.slides.add_slide(blank_layout)
            if ending_bg:
                _apply_background(ending, prs, ending_bg)
            else:
                draw_svg_visuals(ending, prs, strip_svg_tokens(ending_svg))
            thank = "谢谢聆听"
            if "THANK_YOU" in ending_styles:
                style, (vw, vh) = ending_styles["THANK_YOU"]
                x = int(style.x / vw * slide_w)
                y = int(style.y / vh * slide_h)
                font_pt = float(style.font_size_px) * 0.75
                width = int(slide_w * 0.86)
                height = int(Pt(font_pt * 1.6))
                left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                top = max(0, int(y - int(Pt(font_pt * 1.1))))
                _add_styled_textbox(
                    ending,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    lines=[thank],
                    font_name=_pick_font_name(style.font_family),
                    font_size_pt=font_pt,
                    color_hex=style.fill,
                    bold=style.bold,
                    align=_align_from_anchor(style.anchor),
                )
            tagline = str(getattr(data, "subtitle", "") or "").strip()
            if tagline and "TAGLINE" in ending_styles:
                style, (vw, vh) = ending_styles["TAGLINE"]
                x = int(style.x / vw * slide_w)
                y = int(style.y / vh * slide_h)
                font_pt = float(style.font_size_px) * 0.75
                width = int(slide_w * 0.86)
                height = int(Pt(font_pt * 1.4))
                left = int(x - width / 2) if style.anchor == "middle" else int(x - width) if style.anchor == "end" else int(x)
                top = max(0, int(y - int(Pt(font_pt * 1.0))))
                _add_styled_textbox(
                    ending,
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    lines=[tagline],
                    font_name=_pick_font_name(style.font_family),
                    font_size_pt=font_pt,
                    color_hex=style.fill,
                    bold=style.bold,
                    align=_align_from_anchor(style.anchor),
                )

            filename = f"{uuid.uuid4()}.pptx"
            file_path = os.path.join(self.OUTPUT_DIR, filename)
            prs.save(file_path)
            return file_path

        raise ValueError("layout is required")


