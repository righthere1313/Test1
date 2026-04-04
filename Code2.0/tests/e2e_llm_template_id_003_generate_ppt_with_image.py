import argparse
import json
import random
import sys
import time
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from fastapi.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402


def _pick_test_jpg() -> Path:
    tests_dir = Path(__file__).resolve().parent
    jpgs = sorted([p for p in tests_dir.glob("*.jpg") if p.is_file() and p.stat().st_size > 0])
    if jpgs:
        return random.choice(jpgs)
    try:
        from PIL import Image

        p = tests_dir / "e2e_id003.jpg"
        Image.new("RGB", (1400, 900), (235, 220, 205)).save(p, quality=90)
        return p
    except Exception:
        raise RuntimeError("tests 目录没有 jpg，且无法生成示例 jpg（缺少 Pillow）")


def _rect(shape):
    return int(shape.left), int(shape.top), int(shape.width), int(shape.height)


def _overlap(a, b) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def _non_background_pictures(slide, prs: Presentation):
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    pics = []
    for shp in slide.shapes:
        if getattr(shp, "shape_type", None) != 13:
            continue
        try:
            if int(shp.left) == 0 and int(shp.top) == 0 and int(shp.width) == sw and int(shp.height) == sh:
                continue
        except Exception:
            pass
        pics.append(shp)
    return pics


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--topic", default="勾股定理")
    ap.add_argument("--slides_total", type=int, default=10)
    ap.add_argument("--out", default="")
    args = ap.parse_args()

    run_id = str(int(time.time()))
    slides_total = int(args.slides_total)
    expected = max(2, slides_total) - 1

    template_id = "003"

    client = TestClient(app)

    health = client.get("/health")
    if health.status_code != 200:
        raise SystemExit(f"health failed: {health.status_code} {health.text}")

    tpl_resp = client.get("/api/v1/templates/ppt")
    if tpl_resp.status_code != 200:
        raise SystemExit(f"templates failed: {tpl_resp.status_code} {tpl_resp.text}")
    templates = tpl_resp.json()
    if not templates:
        raise SystemExit("templates is empty")

    picked = next((t for t in templates if str(t.get("id")) == template_id), None)
    if not picked:
        raise SystemExit(f"template_id not found: {template_id}")
    template_name = picked.get("file_name") or picked.get("relative_path") or picked.get("name")

    layouts_resp = client.get("/api/v1/templates/ppt/layouts", params={"template": template_id})
    if layouts_resp.status_code != 200:
        raise SystemExit(f"layouts failed: {layouts_resp.status_code} {layouts_resp.text}")
    layouts = layouts_resp.json()

    outline_query = (
        f"请为主题“{args.topic}”生成一个用于{slides_total}页PPT的教学大纲（含封面）。"
        "要求每页给出：页标题、教学目标、推导/例题/活动建议。请输出结构化中文。"
    )
    outline_resp = client.post(
        "/api/v1/chat/qa",
        json={
            "query": outline_query,
            "top_k": 3,
            "document_id": None,
            "temporary_document_ids": None,
            "session_id": f"e2e-outline-id003-{run_id}",
        },
    )
    if outline_resp.status_code != 200:
        raise SystemExit(f"qa outline failed: {outline_resp.status_code} {outline_resp.text}")
    outline = (outline_resp.json() or {}).get("answer") or ""

    plan_req = {
        "template": template_id,
        "topic": args.topic,
        "slides_total": slides_total,
        "extra_instructions": (
            "严格只输出JSON。\n"
            f"slides长度必须是{expected}。\n"
            "每页content写6条要点，每条尽量1-2句，尽量长一些用于压力测试排版。\n"
            "每页都设置 background_image，并尽量使用多种 layout_index。\n"
            "中间页将插入一张图片（右图左文），请保证正文要点足够多。\n"
            "下面是模板布局摘要与教学大纲，请你结合它们设计布局与内容。\n"
            f"模板摘要：{json.dumps(layouts, ensure_ascii=False)}\n"
            f"教学大纲：{outline}\n"
        ),
    }
    plan_resp = client.post("/api/v1/ppt/plan", json=plan_req)
    if plan_resp.status_code != 200:
        raise SystemExit(f"plan failed: {plan_resp.status_code} {plan_resp.text}")
    plan_data = plan_resp.json()
    plan = plan_data.get("plan")
    if not isinstance(plan, dict):
        raise SystemExit(f"plan invalid: {plan_data}")
    slides = plan.get("slides")
    if not isinstance(slides, list) or len(slides) != expected:
        raise SystemExit(f"plan slides length invalid: expected={expected} actual={len(slides) if isinstance(slides, list) else None}")

    jpg_path = _pick_test_jpg()
    image_filename = jpg_path.name
    mid_idx = len(slides) // 2
    if isinstance(slides[mid_idx], dict):
        slides[mid_idx]["image_filename"] = image_filename
        slides[mid_idx]["image_position"] = "right"

    files = {
        "payload": (None, json.dumps(plan, ensure_ascii=False), "application/json"),
        "images": (image_filename, jpg_path.read_bytes(), "image/jpeg"),
    }
    gen_resp = client.post("/api/v1/generate/ppt/multipart", files=files)
    if gen_resp.status_code != 200:
        raise SystemExit(f"generate failed: {gen_resp.status_code} {gen_resp.text}")
    gen_data = gen_resp.json()
    download_url = gen_data.get("download_url")
    if not isinstance(download_url, str) or not download_url.startswith("/"):
        raise SystemExit(f"download_url invalid: {gen_data}")

    dl = client.get(download_url)
    if dl.status_code != 200:
        raise SystemExit(f"download failed: {dl.status_code} {dl.text}")

    out = (args.out or "").strip()
    if not out:
        out = str(Path(__file__).resolve().parent / f"E2E-LLM-template003-{args.topic}-{run_id}.pptx")
    out_path = Path(out).resolve()
    out_path.write_bytes(dl.content)

    prs = Presentation(str(out_path))
    if len(prs.slides) != slides_total:
        raise SystemExit(f"slides_total invalid: expected={slides_total} actual={len(prs.slides)} out={out_path}")

    image_slide_index = mid_idx + 1
    pics = _non_background_pictures(prs.slides[image_slide_index], prs)
    if not pics:
        raise SystemExit(f"no inserted image detected: slide={image_slide_index} out={out_path}")

    pic = min(pics, key=lambda s: int(s.width) * int(s.height))
    pic_rect = _rect(pic)
    title_shape = getattr(prs.slides[image_slide_index].shapes, "title", None)
    for shp in prs.slides[image_slide_index].shapes:
        if shp is title_shape:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        if not getattr(shp.text_frame, "text", "").strip():
            continue
        if _overlap(pic_rect, _rect(shp)):
            raise SystemExit(f"text overlaps inserted image: slide={image_slide_index} out={out_path}")

    plan_out = out_path.with_suffix(".plan.json")
    plan_out.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    print("template_id:", template_id)
    print("template_name:", template_name)
    print("mode:", plan_data.get("mode"))
    print("jpg:", str(jpg_path))
    print("ppt:", str(out_path))
    print("plan:", str(plan_out))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

