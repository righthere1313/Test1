import argparse
import json
import random
import sys
import time
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from fastapi.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402


def _pick_test_jpg() -> Path:
    tests_dir = Path(__file__).resolve().parent
    jpgs = sorted([p for p in tests_dir.glob("*.jpg") if p.is_file() and p.stat().st_size > 0])
    if jpgs:
        return random.choice(jpgs)
    raise RuntimeError("tests 目录下未找到任何 .jpg 图片")


def _non_background_pictures(slide, prs: Presentation):
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    pics = []
    for shp in slide.shapes:
        if getattr(shp, "shape_type", None) != 13:
            continue
        try:
            if int(shp.left) == 0 and int(shp.top) == 0 and int(shp.width) == sw and int(shp.height) == sh:
                continue
        except Exception:
            pass
        pics.append(shp)
    return pics


def _rect(shape):
    return int(shape.left), int(shape.top), int(shape.width), int(shape.height)


def _overlap(a, b) -> bool:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def _amplify_text(plan: dict) -> None:
    slides = plan.get("slides")
    if not isinstance(slides, list):
        return
    for i, s in enumerate(slides):
        if not isinstance(s, dict):
            continue
        content = s.get("content")
        if not isinstance(content, list):
            continue
        if len(content) < 5:
            while len(content) < 5:
                content.append("补充要点：请用一两句话解释“为什么这一步成立”，并给出一个可验证的小例子。")
        if len(content) > 6:
            content = content[:6]
        enriched = []
        for j, item in enumerate(content):
            t = str(item or "").strip()
            if len(t) < 45:
                t = (
                    t
                    + "；请写出关键步骤并说明前提条件（例如“必须是直角三角形”），最后用一个数值例子快速验算。"
                )
            if i in {1, 4} and j == 1 and len(t) < 140:
                t = (
                    t
                    + " 进一步思考：若已知 a 与 c，为什么可以用 b=√(c²−a²)？"
                    + " 请讨论根号内必须非负这一约束，并说明物理/几何意义。"
                )
            enriched.append(t)
        s["content"] = enriched


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template_id", default="")
    ap.add_argument("--template_contains", default="")
    ap.add_argument("--topic", default="勾股定理")
    ap.add_argument("--slides_total", type=int, default=10)
    ap.add_argument("--out", default="")
    args = ap.parse_args()

    run_id = str(int(time.time()))
    client = TestClient(app)

    health = client.get("/health")
    if health.status_code != 200:
        raise SystemExit(f"health failed: {health.status_code} {health.text}")

    templates_resp = client.get("/api/v1/templates/ppt")
    if templates_resp.status_code != 200:
        raise SystemExit(f"templates failed: {templates_resp.status_code} {templates_resp.text}")
    templates = templates_resp.json()
    if not templates:
        raise SystemExit("templates is empty")

    template_id = (args.template_id or "").strip()
    if template_id:
        picked = next((t for t in templates if str(t.get("id")) == template_id), None)
        if not picked:
            raise SystemExit(f"template_id not found: {template_id}")
    else:
        needle = (args.template_contains or "").strip()
        picked = None
        if needle:
            for item in templates:
                file_name = str(item.get("file_name", ""))
                rel = str(item.get("relative_path", ""))
                if needle in file_name or needle in rel:
                    picked = item
                    break
        if not picked:
            picked = random.choice(templates)
        template_id = str(picked["id"])

    template_name = picked.get("file_name") or picked.get("relative_path") or picked.get("name")

    slides_total = int(args.slides_total)
    expected = max(2, slides_total) - 1

    extra = (
        "请严格只输出JSON。\n"
        "每页 content 请写 5-6 条要点，每条尽量用 1-2 句话，包含：前提、推导要点、一个具体数值例子与一句反思/提问。\n"
        "请在JSON中为封面与每一页都设置 background_image（由你选择并返回），并尽量利用模板的 layout_index 做出差异化版式。\n"
        "其中某一页将由客户端插入一张图片（右图左文），请保证中间页正文要点足够多，便于测试排版。\n"
    )

    plan_req = {
        "template": template_id,
        "topic": args.topic,
        "slides_total": slides_total,
        "extra_instructions": extra,
    }
    plan_resp = client.post("/api/v1/ppt/plan", json=plan_req)
    if plan_resp.status_code != 200:
        raise SystemExit(f"plan failed: {plan_resp.status_code} {plan_resp.text}")
    plan_data = plan_resp.json()
    plan = plan_data.get("plan")
    if not isinstance(plan, dict):
        raise SystemExit(f"plan invalid: {plan_data}")

    slides = plan.get("slides")
    if not isinstance(slides, list) or len(slides) != expected:
        raise SystemExit(f"slides length invalid: expected={expected} actual={len(slides) if isinstance(slides, list) else None}")

    _amplify_text(plan)

    jpg_path = _pick_test_jpg()
    image_filename = jpg_path.name

    mid_idx = len(plan["slides"]) // 2
    if isinstance(plan["slides"][mid_idx], dict):
        plan["slides"][mid_idx]["image_filename"] = image_filename
        plan["slides"][mid_idx]["image_position"] = "right"

    files = {
        "payload": (None, json.dumps(plan, ensure_ascii=False), "application/json"),
        "images": (image_filename, jpg_path.read_bytes(), "image/jpeg"),
    }
    gen_resp = client.post("/api/v1/generate/ppt/multipart", files=files)
    if gen_resp.status_code != 200:
        raise SystemExit(f"generate failed: {gen_resp.status_code} {gen_resp.text}")
    gen_data = gen_resp.json()
    download_url = gen_data.get("download_url")
    if not isinstance(download_url, str) or not download_url.startswith("/"):
        raise SystemExit(f"download_url invalid: {gen_data}")

    dl = client.get(download_url)
    if dl.status_code != 200:
        raise SystemExit(f"download failed: {dl.status_code} {dl.text}")

    out = (args.out or "").strip()
    if not out:
        out = str(Path(__file__).resolve().parent / f"客户端全流程-复杂文本-{args.topic}-{run_id}.pptx")
    out_path = Path(out).resolve()
    out_path.write_bytes(dl.content)

    prs = Presentation(str(out_path))
    if len(prs.slides) != slides_total:
        raise SystemExit(f"slides_total invalid: expected={slides_total} actual={len(prs.slides)} out={out_path}")

    image_slide_index = mid_idx + 1
    pics = _non_background_pictures(prs.slides[image_slide_index], prs)
    if not pics:
        raise SystemExit(f"no inserted image detected: slide={image_slide_index} out={out_path}")

    pic = min(pics, key=lambda s: int(s.width) * int(s.height))
    pic_rect = _rect(pic)
    title_shape = getattr(prs.slides[image_slide_index].shapes, "title", None)
    for shp in prs.slides[image_slide_index].shapes:
        if shp is title_shape:
            continue
        if not hasattr(shp, "text_frame"):
            continue
        if not getattr(shp.text_frame, "text", "").strip():
            continue
        if _overlap(pic_rect, _rect(shp)):
            raise SystemExit(f"text overlaps inserted image: slide={image_slide_index} out={out_path}")

    plan_out = out_path.with_suffix(".plan.json")
    plan_out.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    print("template_id:", template_id)
    print("template_name:", template_name)
    print("mode:", plan_data.get("mode"))
    print("jpg:", str(jpg_path))
    print("ppt:", str(out_path))
    print("plan:", str(plan_out))
    print("slides:", len(prs.slides))
    print("image_slide_index:", image_slide_index)
    print("non_bg_pictures_on_image_slide:", len(pics))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

