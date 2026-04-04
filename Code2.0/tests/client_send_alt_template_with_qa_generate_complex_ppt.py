import argparse
import json
import random
import re
import sys
import time
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from fastapi.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402


def _pick_test_jpg() -> Path:
    tests_dir = Path(__file__).resolve().parent
    jpgs = sorted([p for p in tests_dir.glob("*.jpg") if p.is_file() and p.stat().st_size > 0])
    if not jpgs:
        raise RuntimeError("tests 目录下未找到任何 .jpg 图片")
    return random.choice(jpgs)


def _ensure_local_bg() -> Path:
    p = Path(__file__).resolve().parent / "qa_alt_bg.png"
    if p.exists() and p.stat().st_size > 0:
        return p
    from PIL import Image

    Image.new("RGB", (1600, 900), (235, 245, 255)).save(p)
    return p


def _extract_json(text: str) -> dict:
    text = (text or "").strip()
    if not text:
        raise ValueError("empty answer")
    try:
        obj = json.loads(text)
        if isinstance(obj, dict):
            return obj
    except Exception:
        pass
    m = re.search(r"\{[\s\S]*\}", text)
    if not m:
        raise ValueError("no json found")
    obj = json.loads(m.group(0))
    if not isinstance(obj, dict):
        raise ValueError("json root not object")
    return obj


def _non_background_pictures(slide, prs: Presentation) -> int:
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    count = 0
    for shp in slide.shapes:
        if getattr(shp, "shape_type", None) != 13:
            continue
        try:
            if int(shp.left) == 0 and int(shp.top) == 0 and int(shp.width) == sw and int(shp.height) == sh:
                continue
        except Exception:
            pass
        count += 1
    return count


def _normalize_plan(plan: dict, template_id: str, topic: str, slides_total: int, bg_path: str) -> dict:
    plan = dict(plan or {})
    plan["template"] = template_id
    plan.setdefault("title", f"{topic}（课件）")
    plan.setdefault("subtitle", "模板学习 + QA链路")
    plan.setdefault("background_image", bg_path)
    slides = plan.get("slides")
    if not isinstance(slides, list):
        slides = []
    expected = max(2, int(slides_total)) - 1
    if len(slides) > expected:
        slides = slides[:expected]
    while len(slides) < expected:
        slides.append(
            {
                "title": f"补充{len(slides)+1}",
                "content": ["要点1", "要点2", "要点3", "要点4", "要点5", "要点6"],
                "notes": "",
                "layout_index": 1,
                "background_image": bg_path,
            }
        )
    for s in slides:
        if isinstance(s, dict):
            s.setdefault("background_image", bg_path)
    plan["slides"] = slides
    return plan


def _amplify_content(plan: dict) -> None:
    slides = plan.get("slides")
    if not isinstance(slides, list):
        return
    for s in slides:
        if not isinstance(s, dict):
            continue
        content = s.get("content")
        if not isinstance(content, list):
            continue
        while len(content) < 6:
            content.append("补充：说明推导前提、关键步骤，并用一个具体数值进行验算，最后给出一句反思问题。")
        if len(content) > 6:
            content[:] = content[:6]
        for i in range(len(content)):
            t = str(content[i] or "").strip()
            if len(t) < 70:
                t += "；请把公式写完整，并解释每个符号代表什么，以及这一步为什么成立。"
            content[i] = t


def _pick_template(templates: list[dict], needle: str) -> dict:
    needle = (needle or "").strip()
    if needle:
        for item in templates:
            file_name = str(item.get("file_name", ""))
            rel = str(item.get("relative_path", ""))
            if needle in file_name or needle in rel:
                return item
    return random.choice(templates)


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template_contains", default="蓝色简约")
    ap.add_argument("--topic", default="勾股定理")
    ap.add_argument("--slides_total", type=int, default=10)
    ap.add_argument("--out", default="")
    args = ap.parse_args()

    run_id = str(int(time.time()))
    client = TestClient(app)

    health = client.get("/health")
    if health.status_code != 200:
        raise SystemExit(f"health failed: {health.status_code} {health.text}")

    templates_resp = client.get("/api/v1/templates/ppt")
    if templates_resp.status_code != 200:
        raise SystemExit(f"templates failed: {templates_resp.status_code} {templates_resp.text}")
    templates = templates_resp.json()
    if not templates:
        raise SystemExit("templates is empty")

    picked = _pick_template(templates, args.template_contains)
    template_id = str(picked["id"])
    template_name = picked.get("file_name") or picked.get("relative_path") or picked.get("name")

    layouts_resp = client.get("/api/v1/templates/ppt/layouts", params={"template": template_id})
    if layouts_resp.status_code != 200:
        raise SystemExit(f"layouts failed: {layouts_resp.status_code} {layouts_resp.text}")
    layouts = layouts_resp.json()
    layouts_text = json.dumps(layouts, ensure_ascii=False)

    bg_path = str(_ensure_local_bg().resolve())
    jpg_path = _pick_test_jpg()
    image_filename = jpg_path.name

    slides_total = int(args.slides_total)
    expected = max(2, slides_total) - 1

    outline_query = (
        f"请为主题“{args.topic}”生成一个用于10页PPT的教学大纲（含封面），"
        "要求每页给出页标题+本页目标+一个例题/活动建议。请输出结构化的中文文本。"
    )
    outline_resp = client.post(
        "/api/v1/chat/qa",
        json={
            "query": outline_query,
            "top_k": 3,
            "document_id": None,
            "temporary_document_ids": None,
            "session_id": f"qa-outline-{run_id}",
        },
    )
    if outline_resp.status_code != 200:
        raise SystemExit(f"qa outline failed: {outline_resp.status_code} {outline_resp.text}")
    outline = (outline_resp.json() or {}).get("answer") or ""

    plan_prompt = (
        "你是PPT课件编排助手。下面给你一个PPT模板的布局摘要（layout列表、页尺寸等），以及一份教学大纲。\n"
        "请你学习这个模板，并基于大纲输出一个用于后端渲染PPT的JSON（只输出JSON，不要任何解释文字）。\n"
        "JSON结构必须是：{title, subtitle, template, background_image, slides}。\n"
        f"slides是数组，长度必须是{expected}。\n"
        "每个slide元素必须包含：title, content(字符串数组), notes, layout_index, background_image。\n"
        "content每页必须有6条要点，每条尽量1-2句，包含：前提、推导要点、一个数值例子、以及一句反思/提问。\n"
        "layout_index只能从模板摘要的 layouts[].layout_index 中选择，尽量做出差异化版式。\n"
        f"背景图请统一填写为下面这个本地路径：{bg_path}\n"
        f"主题：{args.topic}，总页数含封面为{slides_total}。\n"
        f"template字段请填写模板编号：{template_id}\n"
        "模板布局摘要JSON：\n"
        f"{layouts_text}\n"
        "教学大纲：\n"
        f"{outline}\n"
    )

    qa_resp = client.post(
        "/api/v1/chat/qa",
        json={
            "query": plan_prompt,
            "top_k": 3,
            "document_id": None,
            "temporary_document_ids": None,
            "session_id": f"qa-plan-{run_id}",
        },
    )
    if qa_resp.status_code != 200:
        raise SystemExit(f"qa plan failed: {qa_resp.status_code} {qa_resp.text}")
    answer = (qa_resp.json() or {}).get("answer") or ""

    plan = _normalize_plan(_extract_json(answer), template_id=template_id, topic=args.topic, slides_total=slides_total, bg_path=bg_path)
    _amplify_content(plan)

    mid_idx = len(plan["slides"]) // 2
    if isinstance(plan["slides"][mid_idx], dict):
        plan["slides"][mid_idx]["image_filename"] = image_filename
        plan["slides"][mid_idx]["image_position"] = "right"

    files = {
        "payload": (None, json.dumps(plan, ensure_ascii=False), "application/json"),
        "images": (image_filename, jpg_path.read_bytes(), "image/jpeg"),
    }
    gen_resp = client.post("/api/v1/generate/ppt/multipart", files=files)
    if gen_resp.status_code != 200:
        raise SystemExit(f"generate failed: {gen_resp.status_code} {gen_resp.text}")
    gen_data = gen_resp.json()
    download_url = gen_data.get("download_url")
    if not isinstance(download_url, str) or not download_url.startswith("/"):
        raise SystemExit(f"download_url invalid: {gen_data}")

    dl = client.get(download_url)
    if dl.status_code != 200:
        raise SystemExit(f"download failed: {dl.status_code} {dl.text}")

    out = (args.out or "").strip()
    if not out:
        out = str(Path(__file__).resolve().parent / f"客户端-切换模板-含QA-复杂内容-{args.topic}-{run_id}.pptx")
    out_path = Path(out).resolve()
    out_path.write_bytes(dl.content)

    prs = Presentation(str(out_path))
    if len(prs.slides) != slides_total:
        raise SystemExit(f"slides_total invalid: expected={slides_total} actual={len(prs.slides)} out={out_path}")

    image_slide_index = mid_idx + 1
    sw, sh = int(prs.slide_width), int(prs.slide_height)
    bg_pics = 0
    for shp in prs.slides[1].shapes:
        if getattr(shp, "shape_type", None) == 13:
            if int(shp.left) == 0 and int(shp.top) == 0 and int(shp.width) == sw and int(shp.height) == sh:
                bg_pics += 1
    if bg_pics <= 0:
        raise SystemExit(f"no background picture detected on slide 1: out={out_path}")

    if _non_background_pictures(prs.slides[image_slide_index], prs) <= 0:
        raise SystemExit(f"no inserted image detected: slide={image_slide_index} out={out_path}")

    plan_out = out_path.with_suffix(".plan.json")
    plan_out.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    print("template_id:", template_id)
    print("template_name:", template_name)
    print("bg:", bg_path)
    print("jpg:", str(jpg_path))
    print("ppt:", str(out_path))
    print("plan:", str(plan_out))
    print("slides:", len(prs.slides))
    print("image_slide_index:", image_slide_index)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
