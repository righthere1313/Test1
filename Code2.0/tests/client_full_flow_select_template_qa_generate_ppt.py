import argparse
import json
import random
import sys
import time
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

from fastapi.testclient import TestClient  # noqa: E402

from app.main import app  # noqa: E402


def _pick_test_jpg() -> Path:
    tests_dir = Path(__file__).resolve().parent
    jpgs = sorted([p for p in tests_dir.glob("*.jpg") if p.is_file() and p.stat().st_size > 0])
    if jpgs:
        return random.choice(jpgs)
    try:
        from PIL import Image

        p = tests_dir / "client_generated.jpg"
        Image.new("RGB", (1200, 800), (240, 230, 210)).save(p, quality=90)
        return p
    except Exception:
        raise RuntimeError("tests 目录没有 jpg，且无法生成示例 jpg（缺少 Pillow）")


def _non_background_pictures(slide, prs: Presentation) -> int:
    sw = int(prs.slide_width)
    sh = int(prs.slide_height)
    count = 0
    for shp in slide.shapes:
        if getattr(shp, "shape_type", None) != 13:
            continue
        try:
            if int(shp.left) == 0 and int(shp.top) == 0 and int(shp.width) == sw and int(shp.height) == sh:
                continue
        except Exception:
            pass
        count += 1
    return count


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("--template_name_contains", default="")
    ap.add_argument("--topic", default="勾股定理")
    ap.add_argument("--slides_total", type=int, default=10)
    ap.add_argument("--out", default="")
    args = ap.parse_args()

    run_id = str(int(time.time()))
    client = TestClient(app)

    health = client.get("/health")
    if health.status_code != 200:
        raise SystemExit(f"health failed: {health.status_code} {health.text}")

    templates_resp = client.get("/api/v1/templates/ppt")
    if templates_resp.status_code != 200:
        raise SystemExit(f"templates failed: {templates_resp.status_code} {templates_resp.text}")
    templates = templates_resp.json()
    if not templates:
        raise SystemExit("templates is empty")

    needle = (args.template_name_contains or "").strip()
    picked = None
    if needle:
        for item in templates:
            file_name = str(item.get("file_name", ""))
            rel = str(item.get("relative_path", ""))
            if needle in file_name or needle in rel:
                picked = item
                break
    if not picked:
        picked = random.choice(templates)

    template_id = picked["id"]
    template_name = picked.get("file_name") or picked.get("relative_path") or picked.get("name")

    plan_req = {
        "template": template_id,
        "topic": args.topic,
        "slides_total": int(args.slides_total),
        "extra_instructions": "请在JSON中为封面与每页都设置 background_image，并尽量利用模板布局。中间页将插入一张图片（右图左文）。",
    }
    plan_resp = client.post("/api/v1/ppt/plan", json=plan_req)
    if plan_resp.status_code != 200:
        raise SystemExit(f"plan failed: {plan_resp.status_code} {plan_resp.text}")
    plan_data = plan_resp.json()
    plan = plan_data.get("plan")
    if not isinstance(plan, dict):
        raise SystemExit(f"plan invalid: {plan_data}")

    slides = plan.get("slides")
    expected = max(2, int(args.slides_total)) - 1
    if not isinstance(slides, list) or len(slides) != expected:
        raise SystemExit(f"slides length invalid: expected={expected} actual={len(slides) if isinstance(slides, list) else None}")

    jpg_path = _pick_test_jpg()
    image_filename = jpg_path.name

    mid_idx = len(slides) // 2
    if isinstance(slides[mid_idx], dict):
        slides[mid_idx]["image_filename"] = image_filename
        slides[mid_idx]["image_position"] = "right"

    files = {
        "payload": (None, json.dumps(plan, ensure_ascii=False), "application/json"),
        "images": (image_filename, jpg_path.read_bytes(), "image/jpeg"),
    }
    gen_resp = client.post("/api/v1/generate/ppt/multipart", files=files)
    if gen_resp.status_code != 200:
        raise SystemExit(f"generate failed: {gen_resp.status_code} {gen_resp.text}")
    gen_data = gen_resp.json()
    download_url = gen_data.get("download_url")
    if not isinstance(download_url, str) or not download_url.startswith("/"):
        raise SystemExit(f"download_url invalid: {gen_data}")

    dl = client.get(download_url)
    if dl.status_code != 200:
        raise SystemExit(f"download failed: {dl.status_code} {dl.text}")

    out = (args.out or "").strip()
    if not out:
        out = str(Path(__file__).resolve().parent / f"客户端全流程-{args.topic}-{run_id}.pptx")
    out_path = Path(out).resolve()
    out_path.write_bytes(dl.content)

    prs = Presentation(str(out_path))
    if len(prs.slides) != int(args.slides_total):
        raise SystemExit(f"slides_total invalid: expected={args.slides_total} actual={len(prs.slides)} out={out_path}")

    image_slide_index = mid_idx + 1
    pic_count = _non_background_pictures(prs.slides[image_slide_index], prs)
    if pic_count <= 0:
        raise SystemExit(f"no inserted image detected: slide={image_slide_index} out={out_path}")

    plan_out = out_path.with_suffix(".plan.json")
    plan_out.write_text(json.dumps(plan, ensure_ascii=False, indent=2), encoding="utf-8")

    print("template_id:", template_id)
    print("template_name:", template_name)
    print("mode:", plan_data.get("mode"))
    print("jpg:", str(jpg_path))
    print("ppt:", str(out_path))
    print("plan:", str(plan_out))
    print("slides:", len(prs.slides))
    print("image_slide_index:", image_slide_index)
    print("non_bg_pictures_on_image_slide:", pic_count)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
